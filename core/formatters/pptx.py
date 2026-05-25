from .base import BaseFormatter


class PPTXFormatter(BaseFormatter):

    file_extension = ".pptx"
    write_mode = "wb"
    encoding = None  # binary mode

    def format(self, ppt_data):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides = ppt_data.get("slides", [])
        meta = ppt_data.get("meta", {})

        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("type", "content")
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")

            if slide_type == "cover":
                self._add_cover(prs, title, content, meta)
            elif slide_type == "thanks":
                self._add_thanks(prs, title, content)
            else:
                self._add_content_slide(prs, title, content, slide_type)

        # Return as bytes
        import io
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _add_cover(self, prs, title, content, meta):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        left, top = Inches(2), Inches(2.5)
        width, height = Inches(9.3), Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if content:
            left2, top2 = Inches(2), Inches(4.5)
            txBox2 = slide.shapes.add_textbox(left2, top2, Inches(9.3), Inches(1.2))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = content.replace("\n", " | ")
            p2.font.size = Pt(18)
            p2.alignment = PP_ALIGN.CENTER

    def _add_thanks(self, prs, title, content):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.3), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title or "谢谢"
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs, title, content, slide_type):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True

        # Content
        txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for j, line in enumerate(content.split("\n")):
            line = line.strip()
            if not line:
                continue
            if j == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = f"  {line}"
            p2.font.size = Pt(20)
            p2.space_after = Pt(12)

    def save(self, content, output_dir, prefix="ppt"):
        import os
        from datetime import datetime

        os.makedirs(output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{prefix}_{timestamp}{self.file_extension}"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, "wb") as f:
            f.write(content)
        return filepath
